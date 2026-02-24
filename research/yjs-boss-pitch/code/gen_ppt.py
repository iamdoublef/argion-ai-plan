# -*- coding: utf-8 -*-
"""Generate boss-pitch PPT from ppt-script-v3.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Config ---
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)  # deep navy
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x6D, 0x77)  # teal accent
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x00, 0x6D, 0x77)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT_BG = RGBColor(0xF0, 0xF7, 0xF7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def set_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=BODY_COLOR, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, font_size=18, color=BODY_COLOR, bold=False,
             space_before=Pt(6), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.space_before = space_before
    p.alignment = align
    return p


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def make_table(slide, left, top, width, rows_data, col_widths=None):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.45 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Microsoft YaHei"
                if r == 0:
                    paragraph.font.color.rgb = TABLE_HEADER_FG
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = BODY_COLOR
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
    return table


# Will be continued in gen_ppt_slides.py
if __name__ == "__main__":
    from gen_ppt_slides import build_all_slides
    build_all_slides(prs, BLANK_LAYOUT)
    out = os.path.join(os.path.dirname(__file__), "AI-boss-pitch.pptx")
    prs.save(out)
    print(f"PPT saved to: {out}")
