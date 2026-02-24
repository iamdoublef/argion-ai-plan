# -*- coding: utf-8 -*-
"""Slide content builder - part of gen_ppt.py"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Import helpers from gen_ppt
from gen_ppt import (set_bg, add_text, add_para, add_accent_bar, make_table,
                     TITLE_COLOR, BODY_COLOR, ACCENT, LIGHT_GRAY, BG_WHITE)

L_MARGIN = Inches(0.8)
CONTENT_W = Inches(11.5)


def title_slide(slide, title, subtitle=""):
    set_bg(slide)
    add_text(slide, L_MARGIN, Inches(2.5), CONTENT_W, Inches(1.2),
             title, font_size=44, color=TITLE_COLOR, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, L_MARGIN, Inches(3.8), CONTENT_W, Inches(0.6),
                 subtitle, font_size=20, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


def section_title(slide, title):
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(3.0), Inches(0.08), Inches(0.8))
    add_text(slide, Inches(1.1), Inches(2.9), CONTENT_W, Inches(1.0),
             title, font_size=36, color=TITLE_COLOR, bold=True)


def std_slide(slide, title, bullets, note=""):
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             title, font_size=30, color=TITLE_COLOR, bold=True)
    tf = add_text(slide, L_MARGIN, Inches(1.5), CONTENT_W, Inches(5.0),
                  bullets[0] if bullets else "", font_size=18, color=BODY_COLOR)
    for b in bullets[1:]:
        add_para(tf, b, font_size=18, color=BODY_COLOR)
    if note:
        add_text(slide, L_MARGIN, Inches(6.6), CONTENT_W, Inches(0.5),
                 note, font_size=12, color=RGBColor(0xAA, 0xAA, 0xAA))


def build_all_slides(prs, layout):
    build_p1(prs, layout)
    build_p2(prs, layout)
    build_p3(prs, layout)
    build_p4(prs, layout)
    build_p5(prs, layout)
    build_p6(prs, layout)
    build_p7(prs, layout)
    build_p8(prs, layout)
    build_p9(prs, layout)
    build_p10(prs, layout)
    build_p11(prs, layout)
    build_p12(prs, layout)
    build_p13(prs, layout)
    build_p14(prs, layout)
    build_p15(prs, layout)
    build_p16(prs, layout)


# --- P1: Cover ---
def build_p1(prs, layout):
    slide = prs.slides.add_slide(layout)
    title_slide(slide, "AI\u80fd\u5e2e\u516c\u53f8\u505a\u4ec0\u4e48",
                "\u4e9a\u4fca\u6c0f \u00b7 2026\u5e742\u6708")


# --- P2: Conclusion first ---
def build_p2(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u5148\u8bf4\u7ed3\u8bba", [
        "\u4e00\u53e5\u8bdd\uff1aAI\u7684\u7a97\u53e3\u671f\u5f00\u7740\uff0c\u7ade\u54c1\u5e94\u7528\u7a0b\u5ea6\u666e\u904d\u504f\u4f4e\uff0c\u73b0\u5728\u662f\u5207\u5165\u7684\u65f6\u673a",
        "",
        "\u4e09\u4e2a\u4f9d\u636e\uff1a",
        "\u2022 \u7ade\u54c1AI\u5e94\u7528\u7a0b\u5ea6\u666e\u904d\u504f\u4f4e\uff081-2\u661f/5\u661f\uff09\uff0c\u5148\u884c\u8005\u6709\u5dee\u5f02\u5316\u673a\u4f1a",
        "\u2022 \u6280\u672f\u5df2\u7ecf\u6210\u719f\u5230\u201c\u5f00\u7bb1\u5373\u7528\u201d\uff0c\u4e0d\u9700\u8981\u81ea\u5df1\u8bad\u7ec3\u6a21\u578b",
        "\u2022 \u5df2\u7ecf\u7528\u4e9a\u4fca\u6c0f\u7684\u516c\u5f00\u4ea7\u54c1\u8d44\u6599\u9a8c\u8bc1\u4e86\u4e24\u4e2a\u573a\u666f\uff0c\u4e0d\u662f\u7eb8\u4e0a\u8c08\u5175",
    ])


# --- P3: Company strengths ---
def build_p3(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u5236\u9020\u80fd\u529b\u662f\u62a4\u57ce\u6cb3", [
        "\u2022 60+\u6b3e\u4ea7\u54c1\uff0c\u8986\u76d6\u8154\u5f0f\u5c01\u53e3\u673a\u3001\u5916\u62bd\u5f0f\u3001\u771f\u7a7a\u888b\u5168\u54c1\u7c7b",
        "\u2022 100\u540d\u5de5\u7a0b\u5e08\uff0c3\u4e2a\u751f\u4ea7\u57fa\u5730\uff08\u5e7f\u5dde\u3001\u4f5b\u5c71\u3001\u65b0\u52a0\u5761\uff09",
        "\u2022 Metro\u94c2\u91d1\u4f9b\u5e94\u5546\uff0c\u5168\u7403\u4ec53\u5bb6",
        "\u2022 \u53ef\u964d\u89e3\u771f\u7a7a\u888b\u5168\u7403\u9996\u5bb6 OK Compost + BPI \u53cc\u8ba4\u8bc1",
        "\u2022 17\u5e74\u771f\u7a7a\u5c01\u53e3\u8bbe\u5907\u5236\u9020\u7ecf\u9a8c",
        "",
        "\u4e3b\u8425\u4e1a\u52a1\uff1a\u4ee3\u5de5 + \u8017\u6750\uff0c\u5236\u9020\u80fd\u529b\u662f\u6838\u5fc3\u7ade\u4e89\u529b",
    ])


# --- P4: Business opportunities ---
def build_p4(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u4ee3\u5de5\u548c\u8017\u6750\u662f\u4e3b\u8425\u4e1a\u52a1\uff0cAI\u80fd\u8ba9\u5b83\u4eec\u66f4\u8d5a\u94b1", [
        "\u4ee3\u5de5\u4e1a\u52a1",
        "\u2022 \u5ba2\u6237\u83b7\u53d6\uff1a\u9760\u4eba\u8109\u548c\u5c55\u4f1a\uff0c\u6ca1\u6709\u7cfb\u7edf\u5316\u6e20\u9053",
        "\u2022 ODM\u54cd\u5e94\uff1a2\u5468\u51fa2-3\u4e2a\u65b9\u6848\uff0c\u5ba2\u6237\u7b49\u4e0d\u8d77 \u2192 AI\u8f85\u52a9\u8bbe\u8ba1\uff0c3-5\u5929\u51fa8-10\u4e2a\u65b9\u6848",
        "\u2022 \u5236\u9020\u6548\u7387\uff1a\u5907\u6599\u9760\u7ecf\u9a8c\u3001\u6392\u4ea7\u9760\u8001\u5e08\u5085\u3001\u63d2\u5355\u9760\u5173\u7cfb\u2014\u2014AI\u90fd\u6709\u6210\u719f\u89e3\u6cd5",
        "",
        "\u8017\u6750\u4e1a\u52a1",
        "\u2022 \u7528\u6237\u57fa\u7840\uff1aWevac\u8017\u6750 65,600\u6761Amazon\u8bc4\u4ef7\uff0c\u7528\u6237\u57fa\u7840\u5df2\u7ecf\u5728",
        "\u2022 \u590d\u8d2d\u8fd0\u8425\uff1a\u76ee\u524d\u6ca1\u6709\u7cfb\u7edf\u5316\u590d\u8d2d\u8fd0\u8425\uff0c\u7528\u6237\u4e70\u5b8c\u5c31\u8d70",
        "\u2022 \u5ba2\u6237\u751f\u547d\u5468\u671f\uff1a\u4ece\u201c\u5356\u4e00\u6b21\u201d\u5230\u201c\u6301\u7eed\u590d\u8d2d\u201d\uff0c\u8017\u6750\u6bdb\u5229\u66f4\u9ad8",
    ])


# --- P5: Brand gap ---
def build_p5(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u6709\u5236\u9020\u80fd\u529b\u3001\u6709\u7528\u6237\u57fa\u7840\uff0c\u4e0b\u4e00\u6b65\uff1a\u54c1\u724c",
             font_size=30, color=TITLE_COLOR, bold=True)
    data = [
        ["\u54c1\u724c", "Amazon\u8bc4\u4ef7\u6570"],
        ["FoodSaver", "30,000+"],
        ["GERYON\uff08\u767d\u724c\uff09", "30,400+"],
        ["Nesco", "14,500+"],
        ["Vesta\uff08\u6211\u4eec\uff09", "67"],
    ]
    make_table(slide, L_MARGIN, Inches(1.5), Inches(5.5), data,
               col_widths=[Inches(3.0), Inches(2.5)])
    tf = add_text(slide, Inches(6.5), Inches(1.5), Inches(5.5), Inches(4.0),
                  "\u8fd9\u4e0d\u662f\u4ea7\u54c1\u4e0d\u597d\uff0c\u662f\u6ca1\u6709\u7528\u5bf9\u65b9\u6cd5\u8ba9\u6d88\u8d39\u8005\u770b\u5230",
                  font_size=22, color=TITLE_COLOR, bold=True)
    add_para(tf, "", font_size=10)
    add_para(tf, "AI\u80fd\u5728\u54ea\u91cc\u53d1\u529b\uff1a", font_size=18, color=ACCENT, bold=True)
    add_para(tf, "\u2022 \u5185\u5bb9\u8d28\u91cf\uff1a32\u5206\u949f\u51fa\u5b8c\u6574Listing\uff0c3\u79cd\u8bed\u8a0039\u5206\u949f\u641e\u5b9a", font_size=16)
    add_para(tf, "\u2022 \u7ade\u54c1\u60c5\u62a5\uff1a143,000+\u6761\u8bc4\u4ef7\u5206\u6790\uff0c\u627e\u5230\u7ade\u54c1\u81f4\u547d\u4f24", font_size=16)
    add_para(tf, "\u2022 \u54c1\u724c\u5185\u5bb9\u79ef\u7d2f\uff1a\u7cfb\u7edf\u5316\u4ea7\u51fa\u5185\u5bb9\uff0c\u5efa\u7acb\u6d88\u8d39\u8005\u7aef\u5b58\u5728\u611f", font_size=16)


# --- P6: Pain points table ---
def build_p6(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "5\u4e2a\u9ad8\u4e25\u91cd\u5ea6\u95ee\u9898\uff0cAI\u90fd\u80fd\u89e3\u51b3", font_size=30, color=TITLE_COLOR, bold=True)
    data = [
        ["\u75db\u70b9", "\u73b0\u72b6", "AI\u80fd\u505a\u4ec0\u4e48"],
        ["\u5185\u5bb9\u8d28\u91cf\u504f\u5f31", "\u5b98\u7f51\u5185\u5bb9\u662f\u7ed9\u91c7\u8d2d\u5546\u7684", "AI\u751f\u6210\u7b26\u5408\u5e73\u53f0\u89c4\u8303\u7684\u4ea7\u54c1\u6587\u6848"],
        ["\u7ade\u54c1\u60c5\u62a5\u76f2\u533a", "\u4e0d\u77e5\u9053\u7528\u6237\u5728\u9a82\u4ec0\u4e48\u3001\u5938\u4ec0\u4e48", "AI\u6279\u91cf\u5206\u6790143K+\u6761\u8bc4\u4ef7"],
        ["\u8d28\u68c0\u4f9d\u8d56\u4eba\u5de5", "15\u6761\u6ce8\u5851\u7ebf\u4eba\u5de5\u8d28\u68c0\uff0c\u6f0f\u68c0\u5bfc\u81f4\u5ba2\u8bc9", "AI\u89c6\u89c9\u8d28\u68c0\uff0c\u68c0\u51fa\u738795%\u219299.5%"],
        ["ODM\u8bbe\u8ba1\u54cd\u5e94\u6162", "2\u5468\u51fa2-3\u4e2a\u65b9\u6848", "AI\u8f85\u52a9\u8bbe\u8ba1\uff0c3-5\u5929\u51fa8-10\u4e2a\u65b9\u6848"],
        ["\u65b0\u54c1\u65b9\u5411\u9760\u7ecf\u9a8c", "\u6bcf\u5e74\u63a8\u6570\u5341\u6b3e\u65b0\u54c1\uff0c\u65b9\u5411\u9760\u7ecf\u9a8c", "AI\u5206\u6790\u5e02\u573a\u8d8b\u52bf\uff0c\u6570\u636e\u9a71\u52a8\u51b3\u7b56"],
    ]
    make_table(slide, L_MARGIN, Inches(1.5), Inches(11.5), data,
               col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)])


# --- P7: Industry AI trends ---
def build_p7(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u884c\u4e1aAI\u8d8b\u52bf\uff1a\u7ade\u54c1\u666e\u904d\u504f\u4f4e\uff0c\u7a97\u53e3\u671f\u771f\u5b9e\u5b58\u5728",
             font_size=28, color=TITLE_COLOR, bold=True)
    tf = add_text(slide, L_MARGIN, Inches(1.3), CONTENT_W, Inches(0.5),
                  "\u5168\u7403\u771f\u7a7a\u98df\u54c1\u8bbe\u5907\u5e02\u573a\uff1a160\u4ebf\u7f8e\u5143\uff0c\u5e74\u589e\u957f5.2%",
                  font_size=18, color=ACCENT, bold=True)
    data = [
        ["\u54c1\u724c", "AI\u5e94\u7528\u7a0b\u5ea6", "\u4e3b\u8981\u5e94\u7528"],
        ["FoodSaver (Newell Brands)", "\u2605\u2605\u2606\u2606\u2606", "\u4f9b\u5e94\u94fe\u7aef\uff0c\u5185\u5bb9/\u54c1\u724c\u7aef\u51e0\u4e4e\u6ca1\u6709"],
        ["Anova", "\u2605\u2605\u2606\u2606\u2606", "App\u8fde\u63a5\uff0c\u5185\u5bb9\u8425\u9500\u8f83\u5f3a"],
        ["Breville", "\u2605\u2606\u2606\u2606\u2606", "\u4f20\u7edf\u5236\u9020\u601d\u8def"],
        ["\u5c0f\u718a\u7535\u5668", "\u2605\u2605\u2606\u2606\u2606", "\u56fd\u5185\u7535\u5546\u5185\u5bb9\u751f\u6210"],
        ["\u4e9a\u4fca\u6c0f", "\u2605\u2606\u2606\u2606\u2606", "\u51e0\u4e4e\u4e3a\u96f6"],
    ]
    make_table(slide, L_MARGIN, Inches(2.0), Inches(11.0), data,
               col_widths=[Inches(3.5), Inches(2.5), Inches(5.0)])
    add_text(slide, L_MARGIN, Inches(5.8), CONTENT_W, Inches(0.5),
             "\u7ed3\u8bba\uff1a\u4e3b\u8981\u7ade\u54c1AI\u5e94\u7528\u7a0b\u5ea6\u666e\u904d\u504f\u4f4e\uff0c\u5148\u884c\u8005\u6709\u5dee\u5f02\u5316\u7a97\u53e3",
             font_size=20, color=TITLE_COLOR, bold=True)


# --- P8: 12 AI scenarios ---
def build_p8(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u8bc6\u522b\u51fa12\u4e2aAI\u5e94\u7528\u573a\u666f\uff0c\u8986\u76d66\u5927\u4e1a\u52a1\u73af\u8282", [
        "\u7814\u53d1\u8bbe\u8ba1 \u2192 \u751f\u4ea7\u5236\u9020 \u2192 \u8d28\u91cf\u7ba1\u7406 \u2192 \u9500\u552e\u8425\u9500 \u2192 \u4f9b\u5e94\u94fe \u2192 \u5ba2\u6237\u670d\u52a1",
        "",
        "AI\u8f85\u52a9\u8bbe\u8ba1 | \u667a\u80fd\u6392\u4ea7 | \u89c6\u89c9\u8d28\u68c0 | \u5185\u5bb9\u751f\u6210 | \u9700\u6c42\u9884\u6d4b | \u667a\u80fd\u5ba2\u670d",
        "\u65b0\u54c1\u65b9\u5411\u51b3\u7b56 | \u5de5\u827a\u4f18\u5316 | \u9884\u6d4b\u6027\u7ef4\u62a4 | \u7ade\u54c1\u60c5\u62a5 | \u4f9b\u5e94\u5546\u7ba1\u7406 | \u7528\u6237\u6d1e\u5bdf",
        "\u914d\u65b9AI |  |  | \u54c1\u724c\u5185\u5bb9\u77e9\u9635 |  | ",
        "",
        "\u4f18\u5148\u7ea7\u5206\u7c7b\uff1a",
        "\u2605 \u901f\u8d62\uff081-3\u4e2a\u6708\u89c1\u6548\uff09\uff1a\u5185\u5bb9\u751f\u6210\u3001\u7ade\u54c1\u60c5\u62a5\u3001AI\u8f85\u52a9\u8bbe\u8ba1",
        "\u25c6 \u6218\u7565\uff086-18\u4e2a\u6708\uff09\uff1a\u89c6\u89c9\u8d28\u68c0\u3001\u9700\u6c42\u9884\u6d4b\u3001\u667a\u80fd\u6392\u4ea7",
        "\u25cb \u63a2\u7d22\uff0818\u4e2a\u6708+\uff09\uff1aIoT\u667a\u80fd\u4ea7\u54c1\u3001\u9884\u6d4b\u6027\u7ef4\u62a4",
    ])


# --- P9: Priority matrix ---
def build_p9(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u5148\u505a\u4ec0\u4e48\uff1a\u9ad8\u4ef7\u503c \u00d7 \u9ad8\u53ef\u884c\u6027",
             font_size=30, color=TITLE_COLOR, bold=True)
    # Quadrant labels
    add_text(slide, Inches(1.5), Inches(1.5), Inches(5.0), Inches(0.5),
             "\u2191 \u9ad8\u4ef7\u503c", font_size=14, color=RGBColor(0x99, 0x99, 0x99))
    add_text(slide, Inches(8.0), Inches(6.5), Inches(4.0), Inches(0.5),
             "\u9ad8\u53ef\u884c\u6027 \u2192", font_size=14, color=RGBColor(0x99, 0x99, 0x99))
    # Quadrant content
    add_text(slide, Inches(1.5), Inches(2.0), Inches(4.5), Inches(2.0),
             "\u6218\u7565\u9879\u76ee\n\u00b7 \u89c6\u89c9\u8d28\u68c0\n\u00b7 \u9700\u6c42\u9884\u6d4b\n\u00b7 \u667a\u80fd\u6392\u4ea7",
             font_size=18, color=BODY_COLOR)
    add_text(slide, Inches(7.0), Inches(2.0), Inches(4.5), Inches(2.0),
             "\u2605 \u901f\u8d62\u9879\u76ee\n\u2605 \u5185\u5bb9\u751f\u6210\n\u2605 \u7ade\u54c1\u60c5\u62a5\n\u2605 AI\u8f85\u52a9\u8bbe\u8ba1",
             font_size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(1.5), Inches(4.3), Inches(4.5), Inches(1.5),
             "\u63a2\u7d22\u9879\u76ee\n\u00b7 IoT\u667a\u80fd\u4ea7\u54c1\n\u00b7 \u9884\u6d4b\u6027\u7ef4\u62a4",
             font_size=18, color=RGBColor(0x99, 0x99, 0x99))
    add_text(slide, Inches(7.0), Inches(4.3), Inches(4.5), Inches(1.5),
             "\u57fa\u7840\u5efa\u8bbe\n\u00b7 \u6570\u636e\u6cbb\u7406\n\u00b7 ERP\u6574\u5408",
             font_size=18, color=RGBColor(0x99, 0x99, 0x99))
    add_text(slide, L_MARGIN, Inches(6.2), CONTENT_W, Inches(0.8),
             "\u901f\u8d623\u4e2a\u573a\u666f\u7684\u5171\u540c\u7279\u70b9\uff1a\u6280\u672f\u6210\u719f\u3001\u6709\u73b0\u6210\u5de5\u5177\u3001\u4e0d\u4f9d\u8d56\u5927\u89c4\u6a21\u6570\u636e\u6539\u9020\u30011-3\u4e2a\u6708\u5185\u80fd\u770b\u5230\u6548\u679c",
             font_size=16, color=ACCENT, bold=True)


# --- P10: Success cases ---
def build_p10(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u5927\u5382\u9a8c\u8bc1\u8fc7\u7684\u8def\uff0c\u4e2d\u5c0f\u5382\u7167\u7740\u8d70",
             font_size=30, color=TITLE_COLOR, bold=True)
    data = [
        ["\u4f01\u4e1a", "\u89c4\u6a21", "AI\u5e94\u7528", "\u6210\u679c"],
        ["\u7f8e\u7684\u96c6\u56e2", "\u5927\u578b", "\u667a\u80fd\u6392\u4ea7+AI\u8d28\u68c0+\u9884\u6d4b\u6027\u7ef4\u62a4", "\u751f\u4ea7\u6548\u7387+30%\uff0c\u4e0d\u826f\u7387-50%"],
        ["Panasonic", "\u5927\u578b", "APS\u667a\u80fd\u6392\u4ea7", "\u6392\u4ea7\u5de5\u65f6-75%\uff0c\u751f\u4ea7\u63d0\u524d\u671f-50%"],
        ["200-500\u4eba\u6ce8\u5851\u5382", "\u4e2d\u5c0f", "AI\u89c6\u89c9\u8d28\u68c0", "\u68c0\u6d4b\u901f\u5ea6+5-10\u500d\uff0c\u6f0f\u68c0\u73875%\u21920.5%"],
        ["300-800\u4eba\u5c0f\u5bb6\u7535ODM", "\u4e2d\u5c0f", "AI\u8f85\u52a9\u8bbe\u8ba1", "\u8bbe\u8ba1\u5468\u671f2\u5468\u21923-5\u5929"],
        ["200-500\u4eba\u51fa\u53e3\u4f01\u4e1a", "\u4e2d\u5c0f", "AI\u591a\u8bed\u8a00\u5185\u5bb9\u751f\u6210", "\u5185\u5bb9\u6548\u7387+10\u500d\uff0c\u7ffb\u8bd1\u6210\u672c-80%"],
    ]
    make_table(slide, L_MARGIN, Inches(1.5), Inches(11.5), data,
               col_widths=[Inches(3.0), Inches(1.2), Inches(4.0), Inches(3.3)])
    add_text(slide, L_MARGIN, Inches(5.8), CONTENT_W, Inches(0.5),
             "AI\u5728\u5236\u9020\u4e1a\u7684\u5e94\u7528\u6b63\u5728\u4ece\u201c\u5927\u5382\u5b9e\u9a8c\u201d\u5411\u4e2d\u5c0f\u5382\u6e17\u900f\uff0c\u65e9\u52a8\u6bd4\u665a\u52a8\u6709\u4f18\u52bf",
             font_size=18, color=TITLE_COLOR, bold=True)


# --- P11: Verified case 1 - AI content ---
def build_p11(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u5df2\u9a8c\u8bc1\u6848\u4f8b\u2460 \u2014 C10af AI\u751f\u6210 vs \u5b98\u7f51\u539f\u6587",
             font_size=28, color=TITLE_COLOR, bold=True)
    # Title comparison
    tf = add_text(slide, L_MARGIN, Inches(1.4), Inches(5.5), Inches(1.8),
                  "\u6807\u9898\u5bf9\u6bd4\uff1a", font_size=18, color=ACCENT, bold=True)
    add_para(tf, '\u5b98\u7f51\u539f\u6587\uff1a"C10f"\uff084\u4e2a\u5b57\u7b26\uff09', font_size=16, color=RGBColor(0xCC, 0x33, 0x33))
    add_para(tf, 'AI\u751f\u6210\uff1a"Wevac C10 Chamber Vacuum Sealer - 1000W Oil Pump, -29.9 inHg..."\uff08196\u5b57\u7b26\uff0c\u542b6\u4e2a\u9ad8\u4ef7\u503c\u641c\u7d22\u5173\u952e\u8bcd\uff09', font_size=16, color=ACCENT)
    # Score table
    data = [
        ["\u7ef4\u5ea6", "\u6743\u91cd", "\u5b98\u7f51", "AI\u751f\u6210"],
        ["\u4fe1\u606f\u5b8c\u6574\u5ea6", "25%", "40%", "83%"],
        ["\u8bf4\u670d\u529b/\u8f6c\u5316", "25%", "17%", "90%"],
        ["SEO\u641c\u7d22\u53ef\u89c1\u6027", "20%", "12%", "92%"],
        ["\u4e8b\u5b9e\u51c6\u786e\u6027", "15%", "100%", "92%"],
        ["\u54c1\u724c\u8c03\u6027", "10%", "40%", "95%"],
        ["\u5e73\u53f0\u5408\u89c4\u6027", "5%", "0%", "100%"],
        ["\u52a0\u6743\u603b\u5206", "", "35.7%", "89.6%"],
    ]
    make_table(slide, Inches(6.5), Inches(1.4), Inches(5.8), data,
               col_widths=[Inches(2.0), Inches(1.0), Inches(1.2), Inches(1.6)])
    add_text(slide, L_MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
             "\u6548\u7387\uff1a32\u5206\u949f\u751f\u6210\u5b8c\u6574Listing + 39\u5206\u949f\u8ffd\u52a03\u79cd\u8bed\u8a00\uff08\u4e2d/\u7ca4/\u5fb7\uff09",
             font_size=18, color=TITLE_COLOR, bold=True)


# --- P12: Verified case 2 - VoC ---
def build_p12(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u5df2\u9a8c\u8bc1\u6848\u4f8b\u2461 \u2014 6\u54c1\u724c \u00d7 143,000+\u6761\u8bc4\u4ef7\u5206\u6790",
             font_size=28, color=TITLE_COLOR, bold=True)
    data = [
        ["\u6392\u540d", "\u75db\u70b9", "\u603b\u63d0\u53ca\u91cf", "\u6d89\u53ca\u54c1\u724c"],
        ["1", "\u5bc6\u5c01\u6027\u4e0d\u7a33\u5b9a", "~2,444\u6761", "4\u4e2a\u54c1\u724c\u90fd\u6709"],
        ["2", "\u4ea7\u54c1\u5bff\u547d\u77ed", "~760+\u6761", "Nesco\u6700\u4e25\u91cd\uff08490\u6761\u7eaf\u8d1f\u9762\uff09"],
        ["3", "\u5438\u529b\u4e0d\u8db3/\u4e0d\u4e00\u81f4", "~811\u6761", "Anova\u54c1\u63a7\u4e24\u6781\u5206\u5316"],
    ]
    make_table(slide, L_MARGIN, Inches(1.5), Inches(11.0), data,
               col_widths=[Inches(1.0), Inches(3.0), Inches(2.5), Inches(4.5)])
    tf = add_text(slide, L_MARGIN, Inches(3.8), CONTENT_W, Inches(2.5),
                  "\u53d1\u73b0\u7684\u673a\u4f1a\uff1a", font_size=20, color=ACCENT, bold=True)
    add_para(tf, "\u2022 \u5bc6\u5c01\u6027\u662f\u884c\u4e1a\u901a\u75c5 \u2192 Vesta Smart Seal\u4e13\u5229\u7684\u6700\u4f73\u5356\u70b9\uff08\u4ece\u672a\u5728\u8425\u9500\u4e2d\u63d0\u8fc7\uff09", font_size=18)
    add_para(tf, '\u2022 Nesco"3-4\u4e2a\u6708\u5c31\u574f\u4e86"490\u6761\u5dee\u8bc4 \u2192 \u53ef\u76f4\u63a5\u7528\u4e8e\u5bf9\u6bd4\u8425\u9500', font_size=18)
    add_para(tf, "\u2022 V22\u7684185W/-28.3inHg \u2192 \u540c\u4ef7\u4f4d\u65e0\u5bf9\u624b\uff0c\u4f46\u6ca1\u4eba\u77e5\u9053", font_size=18)
    add_text(slide, L_MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
             "\u54c1\u724c\u8bc4\u4ef7\u6570\uff1aWevac\u8017\u6750 65,600\u6761 / FoodSaver 30,000\u6761 / Vesta 67\u6761",
             font_size=16, color=RGBColor(0x99, 0x99, 0x99))


# --- P13: Top quick wins ---
def build_p13(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u901f\u8d62\u573a\u666f\u7684\u6548\u679c", [
        "\u2460 AI\u591a\u8bed\u8a00\u5185\u5bb9\u751f\u6210",
        "   \u73b0\u72b6\uff1a\u5b98\u7f51\u5185\u5bb9\u662fB2B\u8bed\u8a00\uff0c\u591a\u8bed\u8a00\u9760\u5916\u5305",
        "   AI\u4e4b\u540e\uff1a\u5355\u4e2aSKU listing 32\u5206\u949f\uff0c3\u79cd\u8bed\u8a0039\u5206\u949f\uff0c\u8bc4\u520635.7%\u219289.6%",
        "   \u6536\u76ca\uff1a\u7ffb\u8bd1\u6210\u672c\u964d80%+\uff0c\u5185\u5bb9\u751f\u4ea7\u6548\u7387\u63d0\u534710\u500d",
        "",
        "\u2461 AI\u89c6\u89c9\u8d28\u68c0\uff08\u6ce8\u5851\u7ebf\u8bd5\u70b9\uff09",
        "   \u73b0\u72b6\uff1a15\u6761\u6ce8\u5851\u7ebf\u4eba\u5de5\u8d28\u68c0\uff0c\u68c0\u51fa\u7387\u7ea695%",
        "   AI\u4e4b\u540e\uff1a\u68c0\u51fa\u7387>99.5%\uff0c\u901f\u5ea6<0.5\u79d2/\u4ef6\uff0c\u8d28\u68c0\u6548\u7387\u63d0\u534710\u500d",
        "   \u6210\u719f\u4f9b\u5e94\u5546\uff1a\u767e\u5ea6\u667a\u80fd\u4e91\u3001\u6d77\u5eb7\u5a01\u89c6\u3001\u521b\u65b0\u5947\u667a",
        "",
        "\u2462 AI\u8f85\u52a9\u4ea7\u54c1\u8bbe\u8ba1",
        "   \u73b0\u72b6\uff1aODM\u5ba2\u6237\u8981\u6982\u5ff5\u8bbe\u8ba1\uff0c2\u5468\u51fa2-3\u4e2a\u65b9\u6848",
        "   AI\u4e4b\u540e\uff1a3-5\u5929\u51fa8-10\u4e2a\u65b9\u6848\uff0c\u8bbe\u8ba1\u5e08\u4eba\u5747\u4ea7\u51fa\u63d0\u53473\u500d",
        '   \u6536\u76ca\uff1a\u8bbe\u8ba1\u5e08\u4ece"\u753b\u56fe\u5de5"\u53d8\u6210"\u521b\u610f\u5bfc\u6f14"',
    ])


# --- P14: Manufacturing AI - prediction & scheduling ---
def build_p14(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u4ee3\u5de5\u5382\u6700\u6015\u4e09\u4ef6\u4e8b\uff1a\u5907\u4e86\u6599\u5356\u4e0d\u6389\u3001\u63a5\u4e86\u5355\u6392\u4e0d\u51fa\u6765\u3001\u63d2\u4e86\u6025\u5355\u4e71\u4e86\u5168\u5c40", [
        "\u2460 AI\u9700\u6c42\u9884\u6d4b\uff08\u5907\u6599\u95ee\u9898\uff09",
        "   \u6d77\u8fd0lead time 30-45\u5929\uff0c\u5907\u6599\u9760\u7ecf\u9a8c\uff1b\u5b63\u8282\u6027\u6ce2\u52a830-50%",
        "   AI\u4e4b\u540e\uff1aSaaS\u8ba2\u9605\u5373\u7528\uff0c\u9884\u6d4b\u51c6\u786e\u7387\u63d0\u534720-30%\uff0c\u5e93\u5b58\u6210\u672c\u964d15-25%",
        "",
        "\u2461 AI\u667a\u80fd\u6392\u4ea7\uff08\u6392\u5355\u95ee\u9898\uff09",
        "   24\u6761\u4ea7\u7ebf\u00d7\u591a\u89c4\u683c\u00d7\u591a\u8ba4\u8bc1\uff0c\u6362\u578b\u635f\u8017\u5927\uff0c\u6392\u4ea7\u9760\u8001\u5e08\u5085\u7ecf\u9a8c",
        "   \u53c2\u8003\uff1aPanasonic\u5bfc\u5165APS\u540e\uff0c\u6392\u4ea7\u5de5\u65f68\u4eba\u5929\u21922\u4eba\u5929\uff08-75%\uff09\uff0c\u751f\u4ea7\u63d0\u524d\u671f-50%",
        "",
        "\u2462 \u52a8\u6001\u63d2\u5355\uff08\u63d2\u5355\u95ee\u9898\uff09",
        "   \u8001\u5ba2\u6237\u63d2\u6025\u5355 vs \u65b0\u5ba2\u6237\u5e38\u89c4\u5355\uff0c\u4f18\u5148\u7ea7\u600e\u4e48\u6392\uff1f\u9760\u5173\u7cfb\u8fd8\u662f\u9760\u89c4\u5219\uff1f",
        "   AI\u53ef\u4ee5\u79d2\u7ea7\u91cd\u7b97\u6700\u4f18\u65b9\u6848\uff0c\u81ea\u52a8\u7ed9\u51fa\u4f18\u5148\u7ea7\u5efa\u8bae",
    ])


# --- P15: 36-month roadmap ---
def build_p15(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_accent_bar(slide, L_MARGIN, Inches(0.6), Inches(0.06), Inches(0.5))
    add_text(slide, Inches(1.1), Inches(0.5), CONTENT_W, Inches(0.7),
             "\u4e09\u9636\u6bb5\u63a8\u8fdb\uff0c\u6bcf\u9636\u6bb5\u6709\u51b3\u7b56\u70b9",
             font_size=30, color=TITLE_COLOR, bold=True)
    # Phase 1
    tf = add_text(slide, L_MARGIN, Inches(1.5), Inches(3.5), Inches(4.5),
                  "\u9636\u6bb5\u4e00\uff080-6\u4e2a\u6708\uff09", font_size=20, color=ACCENT, bold=True)
    add_para(tf, "\u5feb\u901f\u9a8c\u8bc1", font_size=18, color=TITLE_COLOR, bold=True)
    add_para(tf, "\u2022 \u542f\u52a8\u5185\u5bb9\u751f\u6210 + \u7ade\u54c1\u60c5\u62a5", font_size=16)
    add_para(tf, "\u2022 \u542f\u52a8AI\u8f85\u52a9\u8bbe\u8ba1\u8bd5\u70b9", font_size=16)
    add_para(tf, "\u2022 \u6570\u636e\u8d44\u4ea7\u76d8\u70b9 + AI\u8ba4\u77e5\u57f9\u8bad", font_size=16)
    add_para(tf, "\u76ee\u6807\uff1a\u8dd1\u51fa\u6570\u636e\uff0c\u5efa\u7acb\u4fe1\u5fc3", font_size=14, color=RGBColor(0x99, 0x99, 0x99))
    # Phase 2
    tf2 = add_text(slide, Inches(4.8), Inches(1.5), Inches(3.5), Inches(4.5),
                   "\u9636\u6bb5\u4e8c\uff086-18\u4e2a\u6708\uff09", font_size=20, color=ACCENT, bold=True)
    add_para(tf2, "\u6838\u5fc3\u6df1\u5316", font_size=18, color=TITLE_COLOR, bold=True)
    add_para(tf2, "\u2022 \u89c6\u89c9\u8d28\u68c0\u5168\u9762\u94fa\u5f00", font_size=16)
    add_para(tf2, "\u2022 \u9700\u6c42\u9884\u6d4b\u4e0a\u7ebf", font_size=16)
    add_para(tf2, "\u2022 \u5efa\u8bbe\u6570\u636e\u5e73\u53f0\uff0c\u6253\u901aERP/MES", font_size=16)
    add_para(tf2, "\u76ee\u6807\uff1aAI\u8986\u76d6\u751f\u4ea7\u548c\u4f9b\u5e94\u94fe\u6838\u5fc3\u73af\u8282", font_size=14, color=RGBColor(0x99, 0x99, 0x99))
    # Phase 3
    tf3 = add_text(slide, Inches(8.8), Inches(1.5), Inches(3.5), Inches(4.5),
                   "\u9636\u6bb5\u4e09\uff0818-36\u4e2a\u6708\uff09", font_size=20, color=ACCENT, bold=True)
    add_para(tf3, "\u5168\u9762\u667a\u80fd\u5316", font_size=18, color=TITLE_COLOR, bold=True)
    add_para(tf3, "\u2022 \u667a\u80fd\u6392\u4ea7\u4e0a\u7ebf", font_size=16)
    add_para(tf3, "\u2022 \u9884\u6d4b\u6027\u7ef4\u62a4\u90e8\u7f72", font_size=16)
    add_para(tf3, "\u2022 \u9996\u6b3eIoT+AI\u667a\u80fd\u4ea7\u54c1\u4e0a\u5e02", font_size=16)
    add_para(tf3, '\u76ee\u6807\uff1a\u4ece"\u5236\u9020\u578b\u4f01\u4e1a"\u5230"\u667a\u80fd\u5236\u9020\u578b\u4f01\u4e1a"', font_size=14, color=RGBColor(0x99, 0x99, 0x99))
    add_text(slide, L_MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
             "\u6bcf\u4e2a\u9636\u6bb5\u672b\u8bbe Go/No-Go \u51b3\u7b56\u70b9\uff1a\u6709\u6570\u636e\u624d\u8fdb\u4e0b\u4e00\u9636\u6bb5\uff0c\u4e0d\u76f2\u76ee\u6295\u5165",
             font_size=18, color=TITLE_COLOR, bold=True)


# --- P16: Next steps ---
def build_p16(prs, layout):
    slide = prs.slides.add_slide(layout)
    std_slide(slide, "\u7acb\u5373\u53ef\u4ee5\u542f\u52a8\u76846\u4ef6\u4e8b", [
        "1. \u786e\u8ba4AI\u8f6c\u578b\u65b9\u5411\uff1aCEO\u5c42\u9762\u660e\u786e\u201cAI\u8d4b\u80fd\u201d\u662f\u6218\u7565\u65b9\u5411\uff0c\u4e0d\u662fIT\u9879\u76ee",
        "",
        "2. \u542f\u52a8AI\u5185\u5bb9\u751f\u6210\u8bd5\u70b9\uff1a\u90093-5\u4e2aWevac SKU\uff0c\u7528AI\u751f\u6210listing\uff0c\u4e0a\u7ebf\u505aA/B\u6d4b\u8bd5",
        "",
        "3. \u7ba1\u7406\u5c42AI\u8ba4\u77e5\u57f9\u8bad\uff1a2\u5929\uff0c\u5916\u90e8\u987e\u95ee\u6388\u8bfe\uff0c\u542b\u5236\u9020\u4e1aAI\u6848\u4f8b\u7814\u8ba8",
        "",
        "4. \u6570\u636e\u8d44\u4ea7\u76d8\u70b9\uff1a\u6478\u6e05\u695a\u516c\u53f8\u6709\u54ea\u4e9b\u6570\u636e\u3001\u5728\u54ea\u91cc\u3001\u8d28\u91cf\u600e\u4e48\u6837",
        "",
        "5. AI\u89c6\u89c9\u8d28\u68c0\u4f9b\u5e94\u5546\u8bc4\u4f30\uff1a\u9080\u8bf72-3\u5bb6\u6765\u73b0\u573a\u52d8\u5bdf\uff0c\u505aPOC\u5bf9\u6bd4",
        "",
        "6. \u5f15\u5165AI\u4e13\u9879\u4eba\u624d\uff1aAI\u8f6c\u578b\u4e0d\u662f\u517c\u804c\u80fd\u5e72\u7684\u4e8b\uff0c\u9700\u8981\u4e00\u4e2a\u61c2AI\u5de5\u5177\u7684\u4eba\u4e32\u8054\u5404\u4e1a\u52a1",
    ], note="\u6700\u5c0f\u7684\u5207\u5165\u70b9\u662f\u5185\u5bb9\u751f\u6210\u2014\u2014\u5de5\u5177\u73b0\u6210\u7684\uff0c\u4e00\u5468\u5185\u5c31\u80fd\u8dd1\u8d77\u6765")
